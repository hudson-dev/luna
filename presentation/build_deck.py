"""Generate the Luna MPC guidance-and-control presentation as an editable .pptx.

Content is drawn from guided_parachutes.md and README.md. Every slide is built
from native text boxes, shapes, and tables (no rasterised images), so the deck
stays editable in PowerPoint, Google Slides, Keynote, and Canva.

    python presentation/build_deck.py [-o output.pptx]
"""

from __future__ import annotations

import argparse
import math
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

MARGIN = Inches(0.85)
BODY_W = SLIDE_W - 2 * MARGIN

THEME = {
    "bg": RGBColor(0x0B, 0x1D, 0x33),
    "bg_alt": RGBColor(0x10, 0x27, 0x42),
    "panel": RGBColor(0x16, 0x31, 0x50),
    "accent": RGBColor(0x35, 0xC8, 0xD6),
    "accent2": RGBColor(0xF2, 0xA9, 0x3B),
    "text": RGBColor(0xF2, 0xF6, 0xFA),
    "muted": RGBColor(0x9B, 0xB0, 0xC4),
    "rule": RGBColor(0x27, 0x45, 0x67),
}

FONT_HEAD = "Verdana"
FONT_BODY = "Verdana"
FONT_MONO = "Consolas"


# --------------------------------------------------------------------------- #
# low-level helpers
# --------------------------------------------------------------------------- #
def new_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs: Presentation, bg: RGBColor | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = bg or THEME["bg"]
    return slide


def textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def write(
    tf,
    text,
    *,
    size=18,
    color=None,
    bold=False,
    italic=False,
    font=FONT_BODY,
    align=PP_ALIGN.LEFT,
    space_before=0,
    space_after=0,
    line_spacing=1.25,
    first=False,
):
    para = tf.paragraphs[0] if first else tf.add_paragraph()
    para.alignment = align
    para.space_before = Pt(space_before)
    para.space_after = Pt(space_after)
    para.line_spacing = line_spacing
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    run.font.color.rgb = color or THEME["text"]
    return para


def wrapped_height(text, width_in, size_pt, *, bold=False, line_spacing=1.25):
    """Approximate rendered height in inches. Verdana averages ~0.58em per glyph."""
    char_w = (size_pt * (0.62 if bold else 0.58)) / 72
    per_line = max(1, int(width_in / char_w))
    longest = max((len(part) for part in text.split("\n")), default=1)
    lines = sum(max(1, math.ceil(len(part) / per_line)) for part in text.split("\n"))
    if longest <= per_line and "\n" not in text:
        lines = 1
    return lines * size_pt * line_spacing / 72


def rect(slide, left, top, width, height, color, *, shape=MSO_SHAPE.RECTANGLE, line=None):
    sh = slide.shapes.add_shape(shape, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1.25)
    sh.shadow.inherit = False
    sh.text_frame.word_wrap = True
    return sh


def kicker(slide, label):
    rect(slide, MARGIN, Inches(0.62), Inches(0.16), Inches(0.16), THEME["accent"])
    tf = textbox(slide, MARGIN + Inches(0.34), Inches(0.55), Inches(8), Inches(0.4))
    write(tf, label.upper(), size=11, color=THEME["accent"], bold=True, first=True)


def heading(slide, title, *, top=Inches(1.05), size=32, width=None):
    tf = textbox(slide, MARGIN, top, width or BODY_W, Inches(1.0))
    write(tf, title, size=size, color=THEME["text"], bold=True, font=FONT_HEAD,
          line_spacing=1.1, first=True)
    return tf


def footer(slide, index, total, note=""):
    tf = textbox(slide, MARGIN, SLIDE_H - Inches(0.62), BODY_W - Inches(1.0), Inches(0.3))
    write(tf, note or "Luna · Guided Parafoil Recovery", size=9, color=THEME["muted"], first=True)
    tf2 = textbox(slide, SLIDE_W - MARGIN - Inches(1.0), SLIDE_H - Inches(0.62),
                  Inches(1.0), Inches(0.3))
    write(tf2, f"{index} / {total}", size=9, color=THEME["muted"], align=PP_ALIGN.RIGHT, first=True)


def bullets(slide, items, *, left=MARGIN, top=Inches(2.15), width=BODY_W,
            size=17, gap=0.26, marker=None):
    """items: list of str or (headline, detail) tuples. Rows self-size to their text."""
    marker = marker or THEME["accent"]
    text_left = left + Inches(0.36)
    text_w_in = (width - Inches(0.36)) / Inches(1)
    y = top
    for item in items:
        head, detail = (item, None) if isinstance(item, str) else item
        rect(slide, left, y + Inches(0.12), Inches(0.13), Inches(0.13), marker)
        tf = textbox(slide, text_left, y, width - Inches(0.36), Inches(0.5))
        write(tf, head, size=size, color=THEME["text"], bold=detail is not None, first=True)
        row = wrapped_height(head, text_w_in, size, bold=detail is not None)
        if detail:
            write(tf, detail, size=size - 3, color=THEME["muted"], space_before=3)
            row += wrapped_height(detail, text_w_in, size - 3) + 3 / 72
        y += Inches(row + gap)
    return y


def card_height(title, lines, width_in, *, title_size=16, body_size=12.5):
    """Inches needed by a card, matching the padding used in card()."""
    inner = width_in - 0.75
    h = 0.34 + wrapped_height(title, inner, title_size, bold=True, line_spacing=1.2)
    for line in lines:
        h += 8 / 72 + wrapped_height(line, inner, body_size, line_spacing=1.3)
    return h + 0.34


def card(slide, left, top, width, height, title, lines, *, accent=None, title_size=16,
         body_size=12.5):
    """height=None auto-sizes to the content."""
    accent = accent or THEME["accent"]
    if height is None:
        height = Inches(card_height(title, lines, width / Inches(1),
                                    title_size=title_size, body_size=body_size))
    rect(slide, left, top, width, height, THEME["panel"])
    rect(slide, left, top, Inches(0.055), height, accent)
    tf = textbox(slide, left + Inches(0.42), top + Inches(0.34), width - Inches(0.75),
                 height - Inches(0.6))
    write(tf, title, size=title_size, color=accent, bold=True, font=FONT_HEAD,
          line_spacing=1.2, first=True)
    for line in lines:
        write(tf, line, size=body_size, color=THEME["text"], space_before=8, line_spacing=1.3)
    return height


def table(slide, headers, rows, *, left=MARGIN, top=Inches(2.1), width=BODY_W,
          col_ratios=None, row_h=0.62, head_h=0.5, size=12.5):
    n_rows, n_cols = len(rows) + 1, len(headers)
    height = Inches(head_h + row_h * len(rows))
    shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    tbl = shape.table
    tbl.first_row = False
    tbl.horz_banding = False

    ratios = col_ratios or [1] * n_cols
    total = sum(ratios)
    for i, ratio in enumerate(ratios):
        tbl.columns[i].width = Emu(int(width * ratio / total))
    tbl.rows[0].height = Inches(head_h)
    for r in range(1, n_rows):
        tbl.rows[r].height = Inches(row_h)

    def fill_cell(cell, text, *, bg, color, bold=False, cell_size=size):
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = cell.margin_right = Inches(0.16)
        cell.margin_top = cell.margin_bottom = Inches(0.06)
        tf = cell.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]
        run = para.add_run()
        run.text = text
        run.font.size = Pt(cell_size)
        run.font.bold = bold
        run.font.name = FONT_BODY
        run.font.color.rgb = color

    for c, head in enumerate(headers):
        fill_cell(tbl.cell(0, c), head, bg=THEME["accent"], color=THEME["bg"],
                  bold=True, cell_size=size - 0.5)
    for r, row in enumerate(rows, start=1):
        bg = THEME["panel"] if r % 2 else THEME["bg_alt"]
        for c, value in enumerate(row):
            fill_cell(tbl.cell(r, c), value, bg=bg,
                      color=THEME["text"] if c == 0 else THEME["muted"],
                      bold=c == 0)
    return tbl


# --------------------------------------------------------------------------- #
# slide builders
# --------------------------------------------------------------------------- #
def slide_title(prs):
    slide = blank_slide(prs, THEME["bg"])
    rect(slide, Emu(0), Emu(0), Inches(0.22), SLIDE_H, THEME["accent"])
    rect(slide, Inches(8.1), Emu(0), Inches(5.233), SLIDE_H, THEME["bg_alt"])

    tf = textbox(slide, MARGIN, Inches(1.9), Inches(6.9), Inches(0.4))
    write(tf, "MODEL PREDICTIVE CONTROL · GNC", size=12, color=THEME["accent"],
          bold=True, first=True)

    tf = textbox(slide, MARGIN, Inches(2.45), Inches(6.9), Inches(2.2))
    write(tf, "Guided Parafoil\nRecovery", size=48, color=THEME["text"], bold=True,
          font=FONT_HEAD, line_spacing=1.05, first=True)

    rect(slide, MARGIN, Inches(4.45), Inches(1.5), Inches(0.045), THEME["accent2"])

    tf = textbox(slide, MARGIN, Inches(4.85), Inches(6.6), Inches(1.4))
    write(tf, "An MPC-based guidance and control stack that steers a deployed "
              "parafoil back to the launch pad under real wind — and a simpler "
              "fallback that flies when the optimiser cannot.",
          size=15, color=THEME["muted"], line_spacing=1.4, first=True)

    tf = textbox(slide, MARGIN, Inches(6.35), Inches(6.6), Inches(0.4))
    write(tf, "Project Luna  ·  Rocketry GNC  ·  Technical Review", size=11,
          color=THEME["muted"], first=True)

    stats = [
        ("6-DOF", "simulator + LTV-QP MPC"),
        ("2 tracks", "primary MPC · fallback guidance"),
        ("OSQP", "embedded solver target"),
    ]
    y = Inches(2.35)
    for value, label in stats:
        tf = textbox(slide, Inches(8.95), y, Inches(4.0), Inches(0.5))
        write(tf, value, size=26, color=THEME["accent"], bold=True, font=FONT_HEAD, first=True)
        write(tf, label, size=12, color=THEME["muted"], space_before=4)
        y += Inches(1.15)
    return slide


def slide_mission(prs):
    slide = blank_slide(prs)
    kicker(slide, "Mission")
    heading(slide, "Bring the vehicle home, predictably")
    tf = textbox(slide, MARGIN, Inches(2.0), Inches(6.4), Inches(2.0))
    write(tf, "Develop an autonomous guided parafoil recovery system for a rocketry "
              "team: steer the deployed parafoil back to the launch pad under "
              "real-world wind disturbances.",
          size=17, color=THEME["text"], line_spacing=1.45, first=True)
    write(tf, "The project spans GNC algorithm development, mechanical and actuator "
              "design, and eventual deployment on flight hardware.",
          size=15, color=THEME["muted"], space_before=14, line_spacing=1.45)

    h = card(slide, Inches(7.6), Inches(1.95), Inches(4.9), None,
             "Theoretical rigor",
             ["A simulation-validated MPC-based guidance and control stack."])
    card(slide, Inches(7.6), Inches(1.95) + h + Inches(0.35), Inches(4.9), None,
         "Practical trustworthiness",
         ["A system that behaves predictably on real flight hardware — and fails "
          "gracefully when it does not."], accent=THEME["accent2"])

    tf = textbox(slide, MARGIN, Inches(4.5), Inches(6.4), Inches(1.4))
    write(tf, "These are treated as distinct engineering goals requiring different "
              "design choices — not two labels for the same work.",
          size=14, color=THEME["accent"], italic=True, line_spacing=1.4, first=True)
    return slide


def slide_two_tracks(prs):
    slide = blank_slide(prs)
    kicker(slide, "Strategy")
    heading(slide, "Two parallel tracks, deliberately")
    tf = textbox(slide, MARGIN, Inches(1.95), Inches(10.5), Inches(0.5))
    write(tf, "A higher-fidelity primary approach and a simpler, more robust fallback "
              "are maintained at the same time. Neither blocks the other.",
          size=15, color=THEME["muted"], line_spacing=1.4, first=True)

    left_title, left_lines = "Track A — Primary: MPC", [
        "Nonlinear receding-horizon optimisation over the parafoil model.",
        "Highest performance and the research vehicle for the project.",
        "Carries solver, compute, and tuning risk onto flight hardware.",
        "Validated first in a 6-DOF Python simulator.",
    ]
    right_title, right_lines = "Track B — Fallback: phased waypoint guidance", [
        "Deterministic guidance law, no optimiser in the loop.",
        "The reliable option for flight hardware.",
        "Predictable failure behaviour and trivially auditable.",
        "Developed in parallel, not as a post-hoc backup.",
    ]
    h = Inches(max(card_height(left_title, left_lines, 5.7, title_size=17),
                   card_height(right_title, right_lines, 5.45, title_size=17)))
    card(slide, MARGIN, Inches(2.85), Inches(5.7), h, left_title, left_lines, title_size=17)
    card(slide, Inches(7.05), Inches(2.85), Inches(5.45), h, right_title, right_lines,
         accent=THEME["accent2"], title_size=17)

    tf = textbox(slide, MARGIN, Inches(2.85) + h + Inches(0.3), Inches(11.5), Inches(0.5))
    write(tf, "Rule of thumb: the track that flies is the one you can explain to the "
              "range safety officer.", size=13, color=THEME["accent"], italic=True, first=True)
    return slide


def slide_domains(prs):
    slide = blank_slide(prs)
    kicker(slide, "Scope")
    heading(slide, "Technical domains in play")
    items = [
        ("Model predictive control and trajectory optimisation",
         "Receding-horizon formulation, constraints, solver selection."),
        ("Guidance law design",
         "Phased waypoint guidance as the robust flight-hardware path."),
        ("Wind estimation",
         "Recovering the disturbance the controller is fighting."),
        ("Parafoil flight modelling",
         "Kinematic and dynamic models, up to 6-DOF with apparent mass."),
        ("Actuator and mechanical design",
         "Servo and brake/control-line system, shock-load protection."),
    ]
    bullets(slide, items, top=Inches(2.15), gap=0.3, width=Inches(10.5))
    return slide


def slide_stack(prs):
    slide = blank_slide(prs)
    kicker(slide, "Architecture")
    heading(slide, "The closed loop, end to end")

    stages = [
        ("Sensing", "GPS, baro,\nIMU"),
        ("Estimation", "State +\nwind estimate"),
        ("Guidance", "Target /\nreference path"),
        ("Control", "MPC or\nwaypoint law"),
        ("Actuation", "Brake-line\nservo"),
    ]
    x = MARGIN
    box_w = Inches(2.12)
    gap = Inches(0.28)
    for i, (name, detail) in enumerate(stages):
        accent = THEME["accent"] if i != 3 else THEME["accent2"]
        sh = rect(slide, x, Inches(2.35), box_w, Inches(1.55), THEME["panel"])
        rect(slide, x, Inches(2.35), box_w, Inches(0.055), accent)
        tf = sh.text_frame
        tf.margin_left = tf.margin_right = Inches(0.14)
        tf.margin_top = Inches(0.22)
        write(tf, name, size=15, color=accent, bold=True, font=FONT_HEAD,
              align=PP_ALIGN.CENTER, first=True)
        write(tf, detail, size=11.5, color=THEME["muted"], align=PP_ALIGN.CENTER,
              space_before=6, line_spacing=1.25)
        if i < len(stages) - 1:
            rect(slide, x + box_w + Inches(0.06), Inches(3.0), Inches(0.16), Inches(0.22),
                 THEME["rule"], shape=MSO_SHAPE.ISOSCELES_TRIANGLE).rotation = 90
        x += box_w + gap

    rect(slide, MARGIN, Inches(4.35), Inches(11.63), Inches(0.03), THEME["rule"])
    tf = textbox(slide, MARGIN, Inches(4.6), Inches(11.6), Inches(0.5))
    write(tf, "Feedback closes the loop every timestep — the plan is recomputed, not replayed.",
          size=13, color=THEME["accent"], italic=True, first=True)

    plant = ["Parafoil dynamics integrated with a fixed-step RK4 scheme, driven by a wind model."]
    valid = ["Closed-loop simulation across calm, steady-wind, shear/turbulence, and gust scenarios."]
    h = Inches(max(card_height("Plant model", plant, 5.7),
                   card_height("Validation", valid, 5.45)))
    card(slide, MARGIN, Inches(5.2), Inches(5.7), h, "Plant model", plant)
    card(slide, Inches(7.05), Inches(5.2), Inches(5.45), h, "Validation", valid,
         accent=THEME["accent2"])
    return slide


def slide_why_mpc(prs):
    slide = blank_slide(prs, THEME["bg_alt"])
    rect(slide, Emu(0), Emu(0), Inches(0.22), SLIDE_H, THEME["accent"])
    kicker(slide, "Key insight")
    tf = textbox(slide, MARGIN, Inches(1.7), Inches(11.2), Inches(2.6))
    write(tf, "MPC's practical robustness comes from receding-horizon feedback on "
              "measured state — not from model accuracy.",
          size=34, color=THEME["text"], bold=True, font=FONT_HEAD, line_spacing=1.2, first=True)
    rect(slide, MARGIN, Inches(4.35), Inches(1.5), Inches(0.045), THEME["accent2"])
    tf = textbox(slide, MARGIN, Inches(4.8), Inches(10.8), Inches(1.6))
    write(tf, "This reframes where effort belongs. Chasing model fidelity has "
              "diminishing returns; guaranteeing that the loop re-solves on fresh "
              "state, on time, every timestep, does not. It also explains why a "
              "coarse model plus fast feedback beats a beautiful model solved too "
              "slowly to matter.",
          size=15, color=THEME["muted"], line_spacing=1.5, first=True)
    return slide


def slide_formulation(prs):
    slide = blank_slide(prs)
    kicker(slide, "Current implementation")
    heading(slide, "MPC formulation as it stands")
    items = [
        ("Nonlinear receding-horizon optimal control",
         "Re-solved each timestep from the current measured state."),
        ("MATLAB `fmincon` with SQP as the solver",
         "Accurate, flexible with nonlinear constraints, and compute-hungry."),
        ("Custom fixed-step RK4 prediction model",
         "Deterministic step cost, which matters for embedded budgeting."),
        ("Wind disturbance model in the loop",
         "The controller is evaluated against the disturbance it must reject."),
    ]
    bullets(slide, items, top=Inches(2.1), gap=0.34, width=Inches(6.6))

    card(slide, Inches(7.85), Inches(2.05), Inches(4.65), None,
         "Why this is a development-only choice",
         ["`fmincon`/SQP iterates linearise → solve → step until convergence.",
          "Iteration count is data-dependent, so worst-case runtime is not bounded.",
          "A flight computer needs a hard per-timestep compute budget.",
          "So: excellent as ground truth, unusable as flight code."],
         accent=THEME["accent2"])
    return slide


def slide_wind(prs):
    slide = blank_slide(prs)
    kicker(slide, "Estimation")
    heading(slide, "Wind estimation: the bias killer")
    tf = textbox(slide, MARGIN, Inches(2.0), Inches(11.5), Inches(0.6))
    write(tf, "Active in the system today. The estimate is formed from the difference "
              "between GPS ground velocity and predicted air velocity, then low-pass filtered.",
          size=15, color=THEME["muted"], line_spacing=1.4, first=True)

    steps = [
        ("GPS ground velocity", "Measured inertial velocity of the vehicle."),
        ("− predicted air velocity", "What the model says the vehicle should be doing."),
        ("→ low-pass filter", "Reject gust-scale noise, keep the persistent component."),
        ("→ wind estimate", "Fed back into prediction so the plan stops drifting."),
    ]
    x = MARGIN
    w = Inches(2.8)
    for i, (name, detail) in enumerate(steps):
        accent = THEME["accent"] if i < 3 else THEME["accent2"]
        rect(slide, x, Inches(3.05), w, Inches(1.75), THEME["panel"])
        rect(slide, x, Inches(3.05), w, Inches(0.055), accent)
        tf = textbox(slide, x + Inches(0.3), Inches(3.35), w - Inches(0.6), Inches(1.2))
        write(tf, name, size=13.5, color=accent, bold=True, font=FONT_HEAD,
              line_spacing=1.2, first=True)
        write(tf, detail, size=11.5, color=THEME["muted"], space_before=8, line_spacing=1.3)
        x += w + Inches(0.12)

    tf = textbox(slide, MARGIN, Inches(5.35), Inches(11.5), Inches(1.0))
    write(tf, "The filter is the design decision. Too fast and the controller chases "
              "turbulence it cannot beat; too slow and a real wind shift is absorbed as "
              "steady-state error. Wind is the dominant disturbance, so this single "
              "estimator removes the largest source of persistent miss distance.",
          size=14, color=THEME["text"], line_spacing=1.45, first=True)
    return slide


def slide_failure_modes(prs):
    slide = blank_slide(prs)
    kicker(slide, "Failure analysis")
    heading(slide, "Three MPC failure modes — and the fix for each")
    table(
        slide,
        ["Failure mode", "What goes wrong", "Mitigation"],
        [
            ("Persistent bias",
             "Unmodelled steady wind pushes the vehicle off target every step.",
             "Wind estimator feeding the prediction model."),
            ("Constraint violation\nunder disturbance",
             "A nominally feasible plan becomes infeasible once disturbed.",
             "Constraint margin, or tube MPC for a guaranteed envelope."),
            ("Terminal-phase\nmodel reliance",
             "Near the ground, model error matters most and horizon shrinks.",
             "Sensor-triggered flare instead of model-based prediction."),
        ],
        top=Inches(2.15), col_ratios=[1.05, 1.65, 1.5], row_h=1.02, head_h=0.5,
    )
    tf = textbox(slide, MARGIN, Inches(5.9), Inches(11.5), Inches(0.9))
    write(tf, "Each mitigation trades optimality for predictability — and each one is "
              "cheap compared with the failure it prevents.",
          size=14, color=THEME["accent"], italic=True, line_spacing=1.4, first=True)
    return slide


def slide_solver_tradeoff(prs):
    slide = blank_slide(prs)
    kicker(slide, "Core tradeoff")
    heading(slide, "SQP versus LTV-QP for flight")
    table(
        slide,
        ["", "SQP (`fmincon`)", "LTV-QP + OSQP"],
        [
            ("Method", "Iterate linearise → solve → step until convergence.",
             "Linearise once per timestep along a reference, solve one QP."),
            ("Accuracy", "Higher — handles the nonlinearity directly.",
             "Lower, but sufficient with a good reference trajectory."),
            ("Compute", "Data-dependent iteration count; no hard bound.",
             "One bounded QP solve per timestep."),
            ("Verdict", "Development and ground-truth reference.",
             "Preferred embedded formulation for flight hardware."),
        ],
        top=Inches(2.05), col_ratios=[0.7, 1.55, 1.75], row_h=0.8, head_h=0.48,
    )
    tf = textbox(slide, MARGIN, Inches(5.95), Inches(11.5), Inches(0.9))
    write(tf, "The migration is driven by compute and reliability constraints, not by "
              "control performance. Bounded worst-case runtime is itself a safety property.",
          size=14, color=THEME["accent"], italic=True, line_spacing=1.4, first=True)
    return slide


def slide_fallback(prs):
    slide = blank_slide(prs)
    kicker(slide, "Track B")
    heading(slide, "Phased waypoint guidance: the option that flies")
    items = [
        ("No optimiser in the flight loop",
         "Behaviour is a direct function of state — inspectable line by line."),
        ("Phased structure",
         "Distinct approach phases replace a single global optimisation."),
        ("Developed in parallel, on purpose",
         "It is the intended flight-hardware path, not a contingency bolted on late."),
        ("Shares the same estimator and actuator stack",
         "Wind estimate and brake-line control are common to both tracks."),
    ]
    bullets(slide, items, top=Inches(2.15), gap=0.36, width=Inches(6.7),
            marker=THEME["accent2"])
    card(slide, Inches(7.85), Inches(2.05), Inches(4.65), None,
         "What we give up — and get",
         ["Give up: constraint handling, optimality, elegant terminal behaviour.",
          "Get: bounded compute, predictable failure, and a control law a reviewer "
          "can reason about without running a solver.",
          "For a recovery system, the second column wins."],
         accent=THEME["accent2"])
    return slide


def slide_mechanical(prs):
    slide = blank_slide(prs)
    kicker(slide, "Mechanical")
    heading(slide, "Safety lives in the load path")
    principle = ["The servo lives only on the brake/control-line branch — never in the "
                 "primary riser load path.",
                 "Flight loads cannot reach the actuator."]
    open_work = ["Servo shock-load protection for the control-line actuator.",
                 "Locking spool mechanism for the brake/control-line system."]
    lesson = ["Self-energising pawl geometry is easy to get wrong — the seating moment "
              "direction is subtle and an incorrect one locks the wrong way.",
              "Of the locking variants, only the friction brake fails gracefully under "
              "shock: it slips rather than transmitting load into the servo.",
              "Graceful failure was chosen over higher holding force."]

    top = Inches(2.05)
    h1 = card(slide, MARGIN, top, Inches(5.7), None, "Core principle", principle,
              title_size=17)
    h2 = card(slide, MARGIN, top + h1 + Inches(0.3), Inches(5.7), None, "Open work",
              open_work, accent=THEME["accent2"], title_size=17)
    right_h = Inches(max(card_height("Hard-won lesson: locking geometry", lesson, 5.45,
                                     title_size=17),
                         (h1 + h2 + Inches(0.3)) / Inches(1)))
    card(slide, Inches(7.05), top, Inches(5.45), right_h,
         "Hard-won lesson: locking geometry", lesson, title_size=17)
    return slide


def slide_roadmap(prs):
    slide = blank_slide(prs)
    kicker(slide, "Roadmap")
    heading(slide, "What comes next, in priority order")
    items = [
        ("1 · Solver migration",
         "Move MPC from `fmincon`/SQP to an LTV-QP formulation solved with OSQP."),
        ("2 · Flare manoeuvre",
         "Rangefinder or barometric trigger, chosen over model-based terminal prediction."),
        ("3 · Successive convexification",
         "Possible middle ground between SQP accuracy and LTV-QP compute cost."),
        ("4 · Model upgrades",
         "Shear-aware wind prediction; 6-DOF dynamics with apparent-mass terms."),
        ("5 · Case study",
         "MIT Underactuated Ch. 10 perching glider — closest analogue to parafoil landing."),
    ]
    bullets(slide, items, top=Inches(2.15), gap=0.28, width=Inches(10.8))
    return slide


def slide_tools(prs):
    slide = blank_slide(prs)
    kicker(slide, "Tooling")
    heading(slide, "Tools and code map")
    table(
        slide,
        ["Tool", "Role"],
        [
            ("MATLAB", "Primary development environment: `fmincon`, Optimization Toolbox, custom RK4 integrator."),
            ("OSQP", "Target QP solver for flight-hardware deployment."),
            ("Python 6-DOF sim", "High-fidelity simulator with LTV-QP MPC, no MATLAB dependency."),
            ("MIT Underactuated Robotics", "Selective theoretical reference, mined rather than followed."),
        ],
        top=Inches(2.05), width=Inches(6.6), col_ratios=[0.8, 1.9], row_h=0.86, head_h=0.46,
        size=11.5,
    )
    files = [
        ("run_parafoil_mpc.m", "top-level MPC run script"),
        ("simulate_flight.m", "closed-loop flight simulation"),
        ("mpc_parafoil.m", "MPC formulation and solver call"),
        ("parafoil_dynamics.m", "vehicle dynamics model"),
        ("wind_model.m", "wind disturbance model"),
        ("rk4.m", "fixed-step RK4 integrator"),
    ]
    top, height = Inches(2.05), Inches(4.15)
    rect(slide, Inches(7.85), top, Inches(4.65), height, THEME["panel"])
    rect(slide, Inches(7.85), top, Inches(0.055), height, THEME["accent2"])
    tf = textbox(slide, Inches(8.27), top + Inches(0.32), Inches(3.9), Inches(0.4))
    write(tf, "Key MATLAB files", size=16, color=THEME["accent2"], bold=True,
          font=FONT_HEAD, first=True)
    y = top + Inches(0.85)
    for name, role in files:
        tf = textbox(slide, Inches(8.27), y, Inches(3.9), Inches(0.44))
        write(tf, name, size=11, color=THEME["text"], font=FONT_MONO, first=True)
        write(tf, role, size=10.5, color=THEME["muted"], space_before=1)
        y += Inches(0.54)
    return slide


def slide_takeaways(prs):
    slide = blank_slide(prs, THEME["bg_alt"])
    rect(slide, Emu(0), Emu(0), SLIDE_W, Inches(0.22), THEME["accent"])
    kicker(slide, "Takeaways")
    heading(slide, "Principles worth carrying forward")
    items = [
        ("Feedback beats fidelity",
         "Re-solving on measured state is the robustness mechanism; model polish is secondary."),
        ("Bound the worst case",
         "A solver without a runtime bound is not a flight solver, however accurate it is."),
        ("Name the failure modes first",
         "Bias, constraint violation, and terminal model reliance each got an explicit mitigation."),
        ("Keep actuators out of the load path",
         "Mechanical safety is architectural, not a matter of stronger parts."),
        ("Fail gracefully by design",
         "The friction brake and the fallback guidance law exist for the same reason."),
    ]
    bullets(slide, items, top=Inches(2.15), gap=0.28, width=Inches(10.8))
    return slide


def slide_closing(prs):
    slide = blank_slide(prs)
    rect(slide, Emu(0), Emu(0), Inches(0.22), SLIDE_H, THEME["accent"])
    tf = textbox(slide, MARGIN, Inches(2.5), Inches(9.5), Inches(1.6))
    write(tf, "Next milestone", size=13, color=THEME["accent"], bold=True, first=True)
    write(tf, "LTV-QP MPC on OSQP, validated in the 6-DOF simulator, with a "
              "sensor-triggered flare and the waypoint fallback ready to fly.",
          size=28, color=THEME["text"], bold=True, font=FONT_HEAD,
          space_before=14, line_spacing=1.25)
    rect(slide, MARGIN, Inches(5.0), Inches(1.5), Inches(0.045), THEME["accent2"])
    tf = textbox(slide, MARGIN, Inches(5.4), Inches(9.5), Inches(0.8))
    write(tf, "Questions, and where you think the risk actually is.",
          size=15, color=THEME["muted"], first=True)
    return slide


BUILDERS = [
    slide_title,
    slide_mission,
    slide_two_tracks,
    slide_domains,
    slide_stack,
    slide_why_mpc,
    slide_formulation,
    slide_wind,
    slide_failure_modes,
    slide_solver_tradeoff,
    slide_fallback,
    slide_mechanical,
    slide_roadmap,
    slide_tools,
    slide_takeaways,
    slide_closing,
]


def build(output: Path) -> Path:
    prs = new_deck()
    total = len(BUILDERS)
    for index, builder in enumerate(BUILDERS, start=1):
        slide = builder(prs)
        if index > 1:
            footer(slide, index, total)
    prs.core_properties.title = "Luna — Guided Parafoil Recovery: MPC Guidance and Control"
    prs.core_properties.subject = "MPC-based guidance and control for autonomous parafoil recovery"
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    return output


def main() -> None:
    default = Path(__file__).resolve().parent / "Luna_MPC_Parafoil_Recovery.pptx"
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("-o", "--output", type=Path, default=default,
                        help="output .pptx path")
    args = parser.parse_args()
    path = build(args.output)
    print(f"wrote {path} ({len(BUILDERS)} slides)")


if __name__ == "__main__":
    main()
